# extracting/msft_extractor.py

# Standard library imports
import io
import tempfile
from pathlib import Path
from typing import List

# Third-party imports
import pandas as pd
import pythoncom
import win32com.client
from docx import Document
import mammoth
from striprtf.striprtf import rtf_to_text

# Local imports
from .basic_extraction import FileTextExtractor
from .extraction_utils import validate_file, normalize_whitespace, run_pandoc, com_app

class WordFileTextExtractor(FileTextExtractor):
    """
    Windows-friendly text extractor for Word formats.
    - DOCX/DOCM: mammoth -> markdown (fallback python-docx)
    - DOC/RTF:   convert via Word COM to TXT (fast, reliable) then read
                 (fallback to pandoc or striprtf if Word isn't installed)
    """
    file_extensions: List[str] = ["docx", "docm", "doc", "rtf"]

    def __init__(self, use_mammoth: bool = True, use_word_com: bool = True,
                 pandoc_path: str | None = None):
        super().__init__()
        self.use_mammoth  = use_mammoth
        self.use_word_com = use_word_com
        self.pandoc_path  = pandoc_path

    def __call__(self, path: str) -> str:
        """
        Determine extraction method for a Word document and return normalized text.

        Parameters
        ----------
        path : str
            Path to the Word (.docx, .docm, .doc) or RTF file.

        Returns
        -------
        str
            Extracted and whitespace-normalized text content.

        Raises
        ------
        FileNotFoundError
            If the input file does not exist.
        ValueError
            If the file extension is unsupported.
        RuntimeError
            If no extraction method succeeds for legacy formats.
        """
        # validate input file
        p = validate_file(path)
        ext = p.suffix.lower().lstrip('.')
        if ext in ("docx", "docm"):
            text = self._extract_docx(str(p))
        elif ext in ("doc", "rtf"):
            text = self._extract_legacy(str(p), ext)
        else:
            raise ValueError(f"Unsupported Word extension: {ext}")
        return normalize_whitespace(text)

    # ---------- helpers ----------
    def _extract_docx(self, path: str) -> str:
        """
        Extract text from a DOCX/DOCM file.

        Tries mammoth to convert to Markdown, falling back to python-docx.

        Parameters
        ----------
        path : str
            Path to the DOCX/DOCM file.

        Returns
        -------
        str
            Raw extracted text or markdown from the document.
        """
        if self.use_mammoth:
            try:
                with open(path, "rb") as f:
                    return mammoth.convert_to_markdown(f).value
            except Exception:
                pass  # fall through to python-docx

        doc = Document(path)
        parts = []
        for p in doc.paragraphs:
            if p.text.strip():
                parts.append(p.text)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    t = cell.text.strip()
                    if t:
                        parts.append(t)
        return "\n".join(parts)

    def _extract_legacy(self, path: str, ext: str) -> str:
        """
        Extract text from legacy Word formats (DOC, RTF).

        RTF uses striprtf; DOC uses COM or pandoc fallback.

        Parameters
        ----------
        path : str
            Path to the legacy file (.doc or .rtf).
        ext : str
            File extension without the dot ('doc' or 'rtf').

        Returns
        -------
        str
            Extracted text content.

        Raises
        ------
        RuntimeError
            If no viable extraction method is available.
        """
        # For RTF files, bypass COM to avoid potential hangs and use striprtf directly
        if ext == "rtf":
            with open(path, "r", encoding="latin-1", errors="ignore") as f:
                return rtf_to_text(f.read())

        if self.use_word_com:
            try:
                return self._word_com_to_txt(path)
            except Exception:
                pass  # fall back

        if self.pandoc_path:
            return self._pandoc_to_txt(path)

        raise RuntimeError("No viable method to extract text from legacy Word file on Windows.")

    def _word_com_to_txt(self, path: str) -> str:
        """
        Use Microsoft Word via COM to SaveAs TXT, then read.

        Parameters
        ----------
        path : str
            Path to the Word document.

        Returns
        -------
        str
            Extracted text content.

        Raises
        ------
        pythoncom.com_error
            If there is an error with COM automation.
        """
        # Constants from Word Object Model (avoid importing win32com.constants each call)
        wdFormatText = 2
        pythoncom.CoInitialize()
        word = win32com.client.DispatchEx("Word.Application")
        word.Visible = False
        txt = ""
        try:
            doc = word.Documents.Open(Path(path).absolute().__str__(), ReadOnly=True)
            with tempfile.TemporaryDirectory() as td:
                out_txt = Path(td) / (Path(path).stem + ".txt")
                doc.SaveAs2(str(out_txt), FileFormat=wdFormatText, Encoding=65001)  # UTF-8
                doc.Close()
                with open(out_txt, "r", encoding="utf-8", errors="ignore") as f:
                    txt = f.read()
        finally:
            word.Quit()
            pythoncom.CoUninitialize()
        return txt

    def _pandoc_to_txt(self, path: str) -> str:
        """
        Convert a document to plain text via Pandoc.

        Parameters
        ----------
        path : str
            Path to the input file to convert.

        Returns
        -------
        str
            Text extracted from the pandoc-converted file.

        Raises
        ------
        subprocess.CalledProcessError
            If the pandoc command fails.
        """
        import subprocess, tempfile
        with tempfile.NamedTemporaryFile(delete=False, suffix=".txt") as tmp:
            out = tmp.name
        cmd = [self.pandoc_path, path, "-t", "plain", "-o", out]
        subprocess.run(cmd, check=True)
        with open(out, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()


class SpreadsheetTextExtractor(FileTextExtractor):
    """
    Flatten spreadsheet content into plain text for embedding.
    """

    file_extensions: List[str] = ["xlsx", "xlsm", "xls", "xlsb", "ods", "csv", "tsv"]

    def __init__(self,
                 sheets: str | List[str] = "all",   # 'all', 'first', or list of names
                 include_headers: bool = True,
                 include_formulas: bool = False,    # needs engine support
                 max_rows: int | None = 5000,
                 max_cols: int | None = 50,
                 delimiter: str = "\t"):
        self.sheets = sheets
        self.include_headers = include_headers
        self.include_formulas = include_formulas
        self.max_rows = max_rows
        self.max_cols = max_cols
        self.delimiter = delimiter

    def __call__(self, path: str) -> str:
        """
        Read and normalize text from spreadsheet or delimited files.

        Parameters
        ----------
        path : str
            Path to the spreadsheet (.xlsx, .xls, .ods, .csv, .tsv, etc.) file.

        Returns
        -------
        str
            Extracted and whitespace-normalized text content.
        """
        # validate input file
        p = validate_file(path)
        ext = p.suffix.lower().lstrip('.')
        if ext in ("csv", "tsv"):
            text = self._read_delimited(p, ext)
        else:
            text = self._read_excel_like(p, ext)
        # normalize whitespace
        return normalize_whitespace(text)

    # ------------- helpers -------------

    def _read_delimited(self, p: Path, ext: str) -> str:
        """
        Read plain text from CSV or TSV files.

        Parameters
        ----------
        p : Path
            Path object to the delimited file.
        ext : str
            Extension without dot ('csv' or 'tsv').

        Returns
        -------
        str
            Raw file content.
        """
        sep = "\t" if ext == "tsv" else ","
        with open(p, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        return text

    def _read_excel_like(self, p: Path, ext: str) -> str:
        """
        Extract text from binary spreadsheet formats via pandas.

        Parameters
        ----------
        p : Path
            Path to the spreadsheet file.
        ext : str
            Extension without dot (e.g., 'xls', 'xlsx', 'ods').

        Returns
        -------
        str
            Combined text content from selected sheets.
        """
        # choose engine
        engine = self._pick_engine(ext)

        # load sheet names
        excel_file = pd.ExcelFile(p, engine=engine)
        sheet_names = excel_file.sheet_names
        if self.sheets == "first":
            sheet_names = sheet_names[:1]
        elif isinstance(self.sheets, list):
            sheet_names = [s for s in sheet_names if s in self.sheets]

        parts = []
        for s in sheet_names:
            df = excel_file.parse(sheet_name=s, engine=engine)

            if self.max_rows: df = df.head(self.max_rows)
            if self.max_cols: df = df.iloc[:, :self.max_cols]

            txt = self._df_to_text(df, sheet=s)
            parts.append(txt)

        return "\n\n".join(parts)

    def _df_to_text(self, df: pd.DataFrame, sheet: str) -> str:
        """
        Serialize a pandas DataFrame to text with optional headers.

        Parameters
        ----------
        df : pandas.DataFrame
            DataFrame containing worksheet data.
        sheet : str
            Name of the sheet being processed.

        Returns
        -------
        str
            Tab-delimited text block with sheet name header.
        """
        # Optionally drop completely empty cols/rows
        df = df.dropna(how="all").dropna(axis=1, how="all")

        buf = io.StringIO()
        buf.write(f"=== Sheet: {sheet} ===\n")
        if self.include_headers:
            buf.write(self.delimiter.join(str(c) for c in df.columns) + "\n")

        for _, row in df.iterrows():
            cells = ["" if pd.isna(v) else str(v) for v in row.tolist()]
            buf.write(self.delimiter.join(cells) + "\n")

        return buf.getvalue()

    def _pick_engine(self, ext: str) -> str:
        """
        Decide which pandas engine to use based on extension and what’s installed.

        Parameters
        ----------
        ext : str
            File extension (e.g., 'xlsx', 'xls').

        Returns
        -------
        str
            Engine name for pandas ExcelFile (e.g., 'openpyxl', 'xlrd').

        Raises
        ------
        ImportError
            If the required engine is not installed and cannot be used.
        """
        if ext in ("xlsx", "xlsm"):
            return "openpyxl"
        if ext == "xls":
            # xlrd >=2 dropped xls; need xlrd==1.2 or fallback
            try:
                import xlrd  # noqa
                return "xlrd"
            except ImportError:
                # try COM/LibreOffice conversion here, else raise
                raise ImportError("xlrd 1.2.0 required for .xls, or convert to .xlsx first.")
        if ext == "xlsb":
            try:
                import pyxlsb  # noqa
                return "pyxlsb"
            except ImportError:
                raise ImportError("pyxlsb required for .xlsb, or convert first.")
        if ext == "ods":
            try:
                import odf  # noqa
                return "odf"
            except ImportError:
                raise ImportError("odfpy required for .ods, or convert first.")
        # fallback
        return "openpyxl"


class PresentationTextExtractor(FileTextExtractor):
    """
    Extract text from presentation files (PPTX/PPT/ODP/...) into plain text.

    Strategy:
    - For pptx/pptm/ppsx: python-pptx
    - For ppt/pps/odp:    convert → pptx or txt via COM or LibreOffice, then parse
    """
    file_extensions: List[str] = ["pptx", "pptm", "ppsx", "ppt", "pps", "odp"]

    def __init__(self,
                 include_notes: bool = True,
                 include_master: bool = False,
                 use_com: bool = True,     # Windows PowerPoint COM
                 soffice_path: str | None = None,  # LibreOffice headless
                 pandoc_path: str | None = None):
        self.include_notes = include_notes
        self.include_master = include_master
        self.use_com = use_com
        self.soffice_path = soffice_path
        self.pandoc_path = pandoc_path

    def __call__(self, path: str) -> str:
        """
        Extract and normalize text from presentation files.

        Parameters
        ----------
        path : str
            Path to the presentation file (.pptx, .ppt, .odp, etc.).

        Returns
        -------
        str
            Extracted and whitespace-normalized text content.
        """
        # validate input file
        p = validate_file(path)
        ext = p.suffix.lower().lstrip('.')

        if ext in ("pptx", "pptm", "ppsx"):
            text = self._extract_pptx(str(p))
        else:
            # ppt, pps, odp → convert
            converted = self._convert_to_pptx_or_txt(str(p), ext)
            if converted.suffix.lower() == ".txt":
                text = converted.read_text(encoding="utf-8", errors="ignore")
            else:
                text = self._extract_pptx(str(converted))
        # normalize whitespace
        return normalize_whitespace(text)

    # ---------- pptx path ----------
    def _extract_pptx(self, path: str) -> str:
        """
        Extract text from PPTX/PPTM/PPSX using python-pptx.

        Parameters
        ----------
        path : str
            Path to the .pptx/.pptm/.ppsx file.

        Returns
        -------
        str
            Combined slide and (optionally) master text.
        """
        from pptx import Presentation
        prs = Presentation(path)

        parts = []
        for idx, slide in enumerate(prs.slides, start=1):
            buf = io.StringIO()
            buf.write(f"=== Slide {idx} ===\n")
            # Slide title (if any)
            if slide.shapes.title and slide.shapes.title.text:
                buf.write(slide.shapes.title.text.strip() + "\n")
            # All shapes
            for shape in slide.shapes:
                txt = self._shape_text(shape)
                if txt:
                    buf.write(txt + "\n")
            # Notes
            if self.include_notes and slide.has_notes_slide:
                notes_txt = slide.notes_slide.notes_text_frame.text
                if notes_txt.strip():
                    buf.write("\n--- Notes ---\n")
                    buf.write(notes_txt.strip() + "\n")
            parts.append(buf.getvalue())

        # Master slides (rarely needed)
        if self.include_master:
            parts.append(self._master_text(prs))

        return "\n\n".join(parts)

    def _shape_text(self, shape) -> str:
        """
        Retrieve text content from a slide shape.

        Parameters
        ----------
        shape : pptx.shapes.base.BaseShape
            A shape object from a python-pptx slide.

        Returns
        -------
        str
            Extracted text, multi-line for tables or grouped shapes.
        """
        # text frame
        if hasattr(shape, "text_frame") and shape.text_frame:
            return "\n".join([p.text for p in shape.text_frame.paragraphs if p.text.strip()])
        # table
        if shape.has_table:
            rows = []
            for r in shape.table.rows:
                cells = [c.text.strip() for c in r.cells]
                rows.append("\t".join(cells))
            return "\n".join(rows)
        # grouped shapes recurse
        if shape.shape_type == 6 and hasattr(shape, "shapes"):  # MSO_SHAPE_TYPE.GROUP = 6
            parts = []
            for sh in shape.shapes:
                txt = self._shape_text(sh)
                if txt:
                    parts.append(txt)
            return "\n".join(parts)
        return ""

    def _master_text(self, prs) -> str:
        """
        Extract text from master slides in a presentation.

        Parameters
        ----------
        prs : pptx.presentation.Presentation
            python-pptx Presentation instance.

        Returns
        -------
        str
            Text content from all master slides.
        """
        buf = io.StringIO()
        buf.write("=== Master Slides ===\n")
        for master in prs.slide_masters:
            for shape in master.shapes:
                txt = self._shape_text(shape)
                if txt:
                    buf.write(txt + "\n")
        return buf.getvalue()

    # ---------- conversion path ----------
    def _convert_to_pptx_or_txt(self, path: str, ext: str) -> Path:
        """
        Return a Path to a .pptx or .txt temp file after conversion.
        Attempt COM (Windows), else LibreOffice, else Pandoc (txt).

        Parameters
        ----------
        path : str
            Path to the source presentation file.

        ext : str
            File extension of the source file (e.g., 'ppt', 'pps').

        Returns
        -------
        Path
            Path to the converted .pptx or .txt file.

        Raises
        ------
        RuntimeError
            If conversion fails and no valid output is produced.
        """
        # Try COM first
        if self.use_com and ext in ("ppt", "pps"):
            try:
                return self._ppt_com_to_pptx(path)
            except Exception:
                pass

        # LibreOffice
        if self.soffice_path:
            try:
                return self._libreoffice_convert(path, "pptx")
            except Exception:
                pass

        # Pandoc fallback to plain text
        if self.pandoc_path:
            try:
                return self._pandoc_to_txt(path)
            except Exception:
                pass

        raise RuntimeError(f"Cannot convert {path}. Install PowerPoint, LibreOffice, or Pandoc.")

    def _ppt_com_to_pptx(self, path: str) -> Path:
        """
        Convert legacy PPT/PPS to PPTX via Windows PowerPoint COM.

        Parameters
        ----------
        path : str
            Path to the .ppt/.pps file.

        Returns
        -------
        Path
            Path to the newly created temporary .pptx file.

        Raises
        ------
        Exception
            If the COM automation fails.
        """
        import pythoncom
        import win32com.client
        pythoncom.CoInitialize()
        powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
        powerpoint.Visible = 0
        tempdir = Path(tempfile.mkdtemp())
        out_path = tempdir / (Path(path).stem + ".pptx")
        try:
            pres = powerpoint.Presentations.Open(str(Path(path).absolute()), WithWindow=False)
            pres.SaveAs(str(out_path), 24)  # ppSaveAsOpenXMLPresentation = 24
            pres.Close()
        finally:
            powerpoint.Quit()
            pythoncom.CoUninitialize()
        return out_path

    def _libreoffice_convert(self, src: str, fmt: str) -> Path:
        """
        Convert a file using LibreOffice headless mode.

        Parameters
        ----------
        src : str
            Source file path.
        fmt : str
            Desired output format (e.g., 'pptx').

        Returns
        -------
        Path
            Path to the converted file in a temp directory.

        Raises
        ------
        subprocess.CalledProcessError
            If the LibreOffice command fails.
        """
        import subprocess, tempfile
        outdir = Path(tempfile.mkdtemp())
        cmd = [self.soffice_path, "--headless", "--convert-to", fmt, "--outdir", str(outdir), src]
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        return next(outdir.glob(f"{Path(src).stem}*.{fmt}"))

    def _pandoc_to_txt(self, src: str) -> Path:
        """
        Fallback conversion of a presentation file to plain text via Pandoc.

        Parameters
        ----------
        src : str
            Source file path.

        Returns
        -------
        Path
            Path to the .txt file created by Pandoc.

        Raises
        ------
        subprocess.CalledProcessError
            If the pandoc command fails.
        """
        import subprocess, tempfile
        out = Path(tempfile.mkdtemp()) / (Path(src).stem + ".txt")
        cmd = [self.pandoc_path, src, "-t", "plain", "-o", str(out)]
        subprocess.run(cmd, check=True)
        return out